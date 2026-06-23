#!/usr/bin/env python3
"""Génère slides/SPIRE-Zero-Trust.pptx — deck sobre, schémas natifs, peu de texte.
Raconte l'évolution v1 (mTLS SPIFFE/SPIRE câblé à la main) -> v2 (service mesh
Istio). Le zero-trust est le fil rouge, porté différemment par les deux versions.
L'oral porte le détail ; les slides ne sont que des repères visuels."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.oxml.ns import qn

# --- palette sobre ---
INK   = RGBColor(0x1A, 0x1F, 0x2B)   # texte principal
DIM   = RGBColor(0x6B, 0x72, 0x80)   # texte secondaire
BRAND = RGBColor(0x2F, 0x6F, 0xED)   # bleu accent
OK    = RGBColor(0x1F, 0xA9, 0x55)   # vert
KO    = RGBColor(0xD9, 0x3A, 0x3A)   # rouge
V2    = RGBColor(0x7C, 0x3A, 0xED)   # violet (accent v2 / mesh)
BG    = RGBColor(0xFF, 0xFF, 0xFF)   # fond blanc
SOFT  = RGBColor(0xF2, 0xF4, 0xF8)   # gris très clair (cartes)
LINE  = RGBColor(0xD7, 0xDC, 0xE5)   # bordures

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
W, H = prs.slide_width, prs.slide_height


def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid(); bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    s.shapes._spTree.remove(bg._element)
    s.shapes._spTree.insert(2, bg._element)
    return s


def textbox(s, x, y, w, h, text, size=18, color=INK, bold=False,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri", spacing=1.0):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = spacing
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color; r.font.name = font
    return tb


def box(s, x, y, w, h, label, sub=None, fill=SOFT, border=LINE, ink=INK, accent=None):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = border; sh.line.width = Pt(1)
    sh.shadow.inherit = False
    if accent:
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Emu(int(w)//28), h)
        bar.fill.solid(); bar.fill.fore_color.rgb = accent
        bar.line.fill.background(); bar.shadow.inherit = False
    tf = sh.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    if accent:
        tf.margin_left = Inches(0.25)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = ink; r.font.name = "Calibri"
    if sub:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = sub
        r2.font.size = Pt(10); r2.font.color.rgb = DIM; r2.font.name = "Calibri"
    return sh


def arrow(s, x1, y1, x2, y2, color=DIM, label=None, dashed=False):
    cn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    cn.line.color.rgb = color; cn.line.width = Pt(1.75)
    cn.shadow.inherit = False
    le = cn.line._get_or_add_ln()
    tail = le.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
    le.append(tail)
    if dashed:
        d = le.makeelement(qn('a:prstDash'), {'val': 'dash'}); le.append(d)
    if label:
        mx, my = (x1 + x2)//2, (y1 + y2)//2
        textbox(s, mx - Inches(0.7), my - Inches(0.28), Inches(1.4), Inches(0.3),
                label, size=9, color=color, align=PP_ALIGN.CENTER)
    return cn


def kicker(s, text, color=BRAND):
    textbox(s, Inches(0.7), Inches(0.45), Inches(9), Inches(0.4),
            text.upper(), size=12, color=color, bold=True)


def title(s, text):
    textbox(s, Inches(0.7), Inches(0.8), Inches(12), Inches(0.9),
            text, size=30, color=INK, bold=True)


def tag(s, text, color):
    """Petite étiquette de version en haut à droite (V1 / V2)."""
    t = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(11.7), Inches(0.45),
                           Inches(1.0), Inches(0.45))
    t.fill.solid(); t.fill.fore_color.rgb = color; t.line.fill.background()
    t.shadow.inherit = False
    tf = t.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = BG; r.font.name = "Calibri"


def footer(s, n):
    textbox(s, Inches(11.8), Inches(7.05), Inches(1.3), Inches(0.3),
            f"{n:02d}", size=10, color=DIM, align=PP_ALIGN.RIGHT)


# =====================================================================
# 1 — Titre
# =====================================================================
s = slide()
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.4), Inches(0.22), Inches(2.7))
band.fill.solid(); band.fill.fore_color.rgb = BRAND; band.line.fill.background(); band.shadow.inherit = False
textbox(s, Inches(0.9), Inches(2.35), Inches(11), Inches(0.5),
        "ZERO-TRUST · mTLS · SERVICE MESH · GITOPS", size=14, color=BRAND, bold=True)
textbox(s, Inches(0.9), Inches(2.85), Inches(11.8), Inches(1.7),
        "Zero-Trust pour microservices", size=42, color=INK, bold=True, spacing=1.0)
textbox(s, Inches(0.9), Inches(4.35), Inches(11.6), Inches(0.9),
        "Du mTLS câblé à la main (SPIFFE/SPIRE) au service mesh (Istio).",
        size=20, color=V2, bold=True, spacing=1.1)
textbox(s, Inches(0.9), Inches(5.2), Inches(11), Inches(0.6),
        "Chaque service prouve son identité avant de parler à un autre. Aucun secret partagé.",
        size=15, color=DIM, spacing=1.1)

# =====================================================================
# 2 — Le zero-trust : le concept (fil rouge)
# =====================================================================
s = slide(); kicker(s, "Le concept"); title(s, "Le zero-trust en deux principes")
box(s, Inches(0.8), Inches(2.2), Inches(5.7), Inches(1.6),
    "Ne jamais faire confiance au réseau",
    "Être dans le cluster ne prouve rien. Toute communication doit être authentifiée et chiffrée.",
    fill=SOFT, accent=BRAND)
box(s, Inches(6.8), Inches(2.2), Inches(5.7), Inches(1.6),
    "Authentifié ≠ autorisé",
    "Prouver son identité ne donne pas tous les droits. Chaque appel est autorisé explicitement.",
    fill=SOFT, accent=OK)
textbox(s, Inches(0.7), Inches(4.3), Inches(12), Inches(1.0),
        "Ces deux principes ne changent pas. Ce qui change entre v1 et v2,\nc'est QUI les met en œuvre : le code applicatif, ou la plateforme (le mesh).",
        size=17, color=INK, spacing=1.3, bold=True)
textbox(s, Inches(0.7), Inches(5.7), Inches(12), Inches(0.8),
        "Fil rouge de la présentation : le même zero-trust, deux implémentations.",
        size=14, color=DIM, align=PP_ALIGN.CENTER)
footer(s, 2)

# =====================================================================
# 3 — Le problème
# =====================================================================
s = slide(); kicker(s, "Le problème"); title(s, "Le réseau ne suffit pas")
textbox(s, Inches(0.7), Inches(1.9), Inches(11.9), Inches(0.8),
        "Dans un cluster, un service peut en joindre un autre par le réseau.\nMais le réseau ne prouve pas QUI appelle.",
        size=18, color=INK, spacing=1.2)
box(s, Inches(1.0), Inches(3.4), Inches(2.6), Inches(1.0), "Catalog", "compromis", fill=SOFT, accent=KO)
box(s, Inches(8.7), Inches(3.4), Inches(2.6), Inches(1.0), "Payments", "sensible", fill=SOFT, accent=KO)
arrow(s, Inches(3.6), Inches(3.9), Inches(8.7), Inches(3.9), color=KO, label="peut joindre ?")
textbox(s, Inches(1.0), Inches(5.0), Inches(11), Inches(1.0),
        "Si un attaquant compromet Catalog, rien au niveau réseau ne l'empêche\nd'atteindre Payments. Il faut une identité applicative vérifiable.",
        size=15, color=DIM, spacing=1.2)
footer(s, 3)

# =====================================================================
# 4 — Les services
# =====================================================================
s = slide(); kicker(s, "L'application"); title(s, "Une mini-boutique, 5 services")
data = [
    ("Gateway", "porte d'entrée", BRAND, 0.7),
    ("Orders", "commandes", INK, 3.18),
    ("Catalog", "produits", INK, 5.66),
    ("Payments", "paiements", KO, 8.14),
    ("Analytics", "Node.js", DIM, 10.62),
]
for name, sub, acc, x in data:
    box(s, Inches(x), Inches(2.5), Inches(2.3), Inches(1.2), name, sub, accent=acc)
textbox(s, Inches(0.7), Inches(4.3), Inches(12), Inches(1.2),
        "Règle sensible :", size=16, color=INK, bold=True)
box(s, Inches(0.7), Inches(4.85), Inches(11.9), Inches(0.95),
    "Seul Orders peut déclencher un paiement.",
    "Gateway orchestre, Catalog expose — mais aucun ne peut appeler /pay directement.",
    fill=SOFT, accent=OK)
textbox(s, Inches(0.7), Inches(6.1), Inches(12), Inches(0.5),
        "Cette règle, c'est notre test du zero-trust — on la garde identique en v1 et en v2.",
        size=14, color=DIM, align=PP_ALIGN.CENTER)
footer(s, 4)

# =====================================================================
# 5 — V1 — SPIRE : l'identité
# =====================================================================
s = slide(); kicker(s, "Version 1 · identité"); title(s, "SPIRE attribue une identité cryptographique"); tag(s, "V1", BRAND)
box(s, Inches(0.9), Inches(2.4), Inches(2.6), Inches(1.0), "SPIRE Server", "autorité / CA", accent=BRAND)
box(s, Inches(0.9), Inches(4.0), Inches(2.6), Inches(1.0), "SPIRE Agent", "DaemonSet")
arrow(s, Inches(2.2), Inches(3.4), Inches(2.2), Inches(4.0), color=BRAND, label="atteste")
for i, (svc, y) in enumerate([("gateway", 2.4), ("orders", 3.5), ("payments", 4.6)]):
    box(s, Inches(5.3), Inches(y), Inches(2.4), Inches(0.9), svc)
    arrow(s, Inches(3.5), Inches(4.5), Inches(5.3), Inches(y + 0.45), color=DIM)
textbox(s, Inches(8.1), Inches(2.5), Inches(4.5), Inches(2.5),
        "SVID X.509", size=18, color=INK, bold=True)
box(s, Inches(8.1), Inches(3.1), Inches(4.4), Inches(1.5),
    "spiffe://example.org/\nns/shop/sa/orders",
    "identité courte durée, rotative", fill=SOFT, accent=BRAND)
textbox(s, Inches(0.7), Inches(5.7), Inches(12), Inches(1),
        "L'identité dérive du ServiceAccount Kubernetes. L'Agent ne distribue un SVID\nqu'après avoir attesté le nœud (PSAT) et le workload.",
        size=14, color=DIM, spacing=1.2)
footer(s, 5)

# =====================================================================
# 6 — V1 — mTLS + autorisation dans le code
# =====================================================================
s = slide(); kicker(s, "Version 1 · le cœur"); title(s, "mTLS transporte l'identité, le code l'autorise"); tag(s, "V1", BRAND)
box(s, Inches(0.8), Inches(2.4), Inches(2.5), Inches(1.1), "Orders", accent=BRAND)
box(s, Inches(9.9), Inches(2.4), Inches(2.5), Inches(1.1), "Payments", accent=KO)
arrow(s, Inches(3.3), Inches(2.95), Inches(9.9), Inches(2.95), color=OK, label="mTLS  ·  /pay")
textbox(s, Inches(3.4), Inches(3.3), Inches(6.4), Inches(0.5),
        "présente son SVID — Payments lit l'identité dans le certificat",
        size=11, color=DIM, align=PP_ALIGN.CENTER)
box(s, Inches(3.0), Inches(4.2), Inches(7.3), Inches(1.6),
    'var policy = map[string][]string{\n    "/pay": {ordersID},\n}',
    "Payments (en Go) : seule l'identité Orders est autorisée sur /pay", fill=SOFT, accent=OK)
textbox(s, Inches(0.7), Inches(6.1), Inches(12), Inches(0.6),
        "Le zero-trust est porté PAR LE CODE : chaque service ouvre le Workload API, fait son mTLS, applique sa policy.",
        size=14, color=INK, bold=True, align=PP_ALIGN.CENTER)
footer(s, 6)

# =====================================================================
# 7 — V1 — Analytics derrière Envoy manuel
# =====================================================================
s = slide(); kicker(s, "Version 1 · multi-langage"); title(s, "Analytics (Node.js) : un sidecar Envoy à la main"); tag(s, "V1", BRAND)
box(s, Inches(0.8), Inches(3.0), Inches(2.4), Inches(1.1), "Gateway", accent=BRAND)
box(s, Inches(5.2), Inches(2.4), Inches(6.5), Inches(2.4), "", fill=RGBColor(0xF7,0xF9,0xFC), border=LINE)
box(s, Inches(5.6), Inches(2.9), Inches(2.4), Inches(1.0), "Envoy", "SDS · mTLS · 99 lignes YAML", accent=BRAND)
box(s, Inches(8.7), Inches(2.9), Inches(2.6), Inches(1.0), "Analytics", "Node.js", accent=DIM)
arrow(s, Inches(3.2), Inches(3.55), Inches(5.6), Inches(3.4), color=OK, label="mTLS")
arrow(s, Inches(8.0), Inches(3.4), Inches(8.7), Inches(3.4), color=DIM, label="HTTP local")
textbox(s, Inches(0.7), Inches(5.3), Inches(12), Inches(1),
        "Pour ne pas réécrire SPIFFE en Node.js, on met un Envoy devant — mais configuré\nà la main (SDS, contextes TLS). Ça marche, mais c'est lourd à maintenir.",
        size=14, color=DIM, spacing=1.2)
footer(s, 7)

# =====================================================================
# 8 — V1 — les limites (slide charnière)
# =====================================================================
s = slide(); kicker(s, "Pourquoi faire évoluer", color=V2); title(s, "Les limites du tout-manuel"); tag(s, "V1", BRAND)
limits = [
    ("mTLS recâblé partout", "chaque service Go porte le code SPIFFE/TLS"),
    ("Autorisation en dur", "la policy vit dans le binaire — recompiler pour changer"),
    ("Sidecar Envoy manuelle", "99 lignes de YAML fragiles à maintenir"),
    ("Peu observable", "aucune vue d'ensemble du « qui parle à qui »"),
]
y = 2.2
for name, sub in limits:
    box(s, Inches(0.8), Inches(y), Inches(11.7), Inches(0.85), name, sub, fill=RGBColor(0xFC,0xED,0xED), border=KO, ink=KO)
    y += 1.0
textbox(s, Inches(0.7), Inches(6.4), Inches(12), Inches(0.6),
        "Le zero-trust est correct, mais c'est l'application qui porte la sécurité. → On délègue ça à un service mesh.",
        size=14, color=V2, bold=True, align=PP_ALIGN.CENTER)
footer(s, 8)

# =====================================================================
# 9 — V2 — le service mesh Istio (la bascule)
# =====================================================================
s = slide(); kicker(s, "Version 2 · la bascule", color=V2); title(s, "Istio : la plateforme porte la sécurité"); tag(s, "V2", V2)
# avant
box(s, Inches(0.8), Inches(2.3), Inches(5.5), Inches(1.9),
    "Avant — l'app porte tout",
    "Code mTLS dans chaque service · policy en Go · Envoy manuel · SPIRE", fill=SOFT, accent=BRAND)
arrow(s, Inches(6.5), Inches(3.25), Inches(7.4), Inches(3.25), color=V2)
box(s, Inches(7.6), Inches(2.3), Inches(5.0), Inches(1.9),
    "Après — le mesh porte tout",
    "Sidecar injecté automatiquement · mTLS auto · l'app fait du HTTP simple", fill=RGBColor(0xF3,0xEE,0xFD), accent=V2)
textbox(s, Inches(0.7), Inches(4.7), Inches(12), Inches(1.4),
        "Istio injecte un proxy (Envoy) à côté de chaque service. Le proxy fait le mTLS,\nporte l'identité (SPIFFE, dérivée du ServiceAccount) et applique l'autorisation.\nRésultat : −1265 lignes de code applicatif, +223 — l'app redevient pur métier.",
        size=15, color=INK, spacing=1.3)
footer(s, 9)

# =====================================================================
# 10 — V2 — mTLS auto + AuthorizationPolicy déclaratives
# =====================================================================
s = slide(); kicker(s, "Version 2 · le cœur", color=V2); title(s, "Le même zero-trust, en déclaratif"); tag(s, "V2", V2)
box(s, Inches(0.8), Inches(2.3), Inches(2.5), Inches(1.1), "Orders", accent=V2)
box(s, Inches(9.9), Inches(2.3), Inches(2.5), Inches(1.1), "Payments", accent=KO)
arrow(s, Inches(3.3), Inches(2.85), Inches(9.9), Inches(2.85), color=V2, label="mTLS auto (mesh)")
box(s, Inches(2.6), Inches(3.9), Inches(8.1), Inches(1.9),
    "kind: AuthorizationPolicy\nrules: from principals [ sa/orders ]  to paths [ /pay ]",
    "PeerAuthentication STRICT (mTLS imposé) + AuthorizationPolicy (qui peut quoi)", fill=RGBColor(0xF3,0xEE,0xFD), accent=V2)
textbox(s, Inches(0.7), Inches(6.1), Inches(12), Inches(0.6),
        "Même règle qu'en v1 (« /pay réservé à Orders »), mais hors du code : du YAML versionné en Git, auditable, sans recompiler.",
        size=14, color=INK, bold=True, align=PP_ALIGN.CENTER)
footer(s, 10)

# =====================================================================
# 11 — V2 — le 403 : zero-trust démontré
# =====================================================================
s = slide(); kicker(s, "Version 2 · démonstration", color=V2); title(s, "Authentifié, mais pas autorisé"); tag(s, "V2", V2)
box(s, Inches(0.8), Inches(2.4), Inches(11.7), Inches(0.95),
    "Gateway → Orders → Payments  ·  /pay     ✓ autorisé", fill=RGBColor(0xEC,0xF8,0xF0), border=OK, ink=OK)
box(s, Inches(0.8), Inches(3.6), Inches(11.7), Inches(0.95),
    "Gateway → Payments  ·  /pay     ✗ 403 — refusé par le mesh", fill=RGBColor(0xFC,0xED,0xED), border=KO, ink=KO)
textbox(s, Inches(0.8), Inches(4.8), Inches(11.7), Inches(1.2),
        "La Gateway est authentifiée (mTLS OK), mais l'AuthorizationPolicy ne l'autorise pas\nsur /pay. Le proxy renvoie 403 AVANT que la requête n'atteigne l'application.",
        size=15, color=DIM, spacing=1.25)
textbox(s, Inches(0.7), Inches(6.2), Inches(12), Inches(0.5),
        "C'est le mesh — plus le code — qui décide. Même réseau, même cluster : seule l'identité change la décision.",
        size=14, color=INK, bold=True, align=PP_ALIGN.CENTER)
footer(s, 11)

# =====================================================================
# 12 — V2 — Kiali : le zero-trust observable
# =====================================================================
s = slide(); kicker(s, "Version 2 · observabilité", color=V2); title(s, "Kiali : voir le mesh en direct"); tag(s, "V2", V2)
cards = [
    ("Topologie live", "qui parle à qui, en temps réel", V2),
    ("Cadenas mTLS", "chiffrement visible par lien", OK),
    ("403 en rouge", "le refus zero-trust, à l'œil", KO),
    ("Trafic continu", "loadgen : un client réaliste", BRAND),
]
x = 0.8
for name, sub, acc in cards:
    box(s, Inches(x), Inches(2.6), Inches(2.85), Inches(1.5), name, sub, accent=acc)
    x += 3.0
textbox(s, Inches(0.7), Inches(4.6), Inches(12), Inches(1.3),
        "En v1, le « qui parle à qui » était invisible. En v2, Kiali (+ Grafana pour les\ntendances) le rend lisible : la topologie s'anime, le mTLS et les refus sont visibles.\nLe code applicatif, lui, ne produit plus aucune métrique maison.",
        size=15, color=INK, spacing=1.25)
footer(s, 12)

# =====================================================================
# 13 — Le zero-trust : v1 vs v2 (le cœur du narratif)
# =====================================================================
s = slide(); kicker(s, "La comparaison"); title(s, "Le même zero-trust, deux implémentations")
# en-têtes de colonnes
box(s, Inches(0.8), Inches(1.9), Inches(3.9), Inches(0.7), "Aspect", fill=INK, ink=BG, border=INK)
box(s, Inches(4.85), Inches(1.9), Inches(3.9), Inches(0.7), "V1 — manuel (SPIRE)", fill=BRAND, ink=BG, border=BRAND)
box(s, Inches(8.9), Inches(1.9), Inches(3.6), Inches(0.7), "V2 — mesh (Istio)", fill=V2, ink=BG, border=V2)
rows = [
    ("Identité", "SVID SPIRE dans le code", "SPIFFE via le proxy"),
    ("mTLS", "câblé en Go / Envoy", "automatique, injecté"),
    ("Autorisation", "policy en dur (Go)", "AuthorizationPolicy YAML"),
    ("Observabilité", "métriques maison", "Kiali + Grafana"),
]
y = 2.75
for aspect, v1, v2 in rows:
    box(s, Inches(0.8), Inches(y), Inches(3.9), Inches(0.82), aspect, fill=SOFT)
    box(s, Inches(4.85), Inches(y), Inches(3.9), Inches(0.82), v1, fill=BG)
    box(s, Inches(8.9), Inches(y), Inches(3.6), Inches(0.82), v2, fill=RGBColor(0xF3,0xEE,0xFD))
    y += 0.92
textbox(s, Inches(0.7), Inches(6.5), Inches(12), Inches(0.5),
        "Les principes sont identiques. La v2 les sort du code et les confie à la plateforme.",
        size=14, color=DIM, align=PP_ALIGN.CENTER)
footer(s, 13)

# =====================================================================
# 14 — Accès zero-trust : JIT SSH + Cloudflare Tunnel/Access
# =====================================================================
s = slide(); kicker(s, "Zero-trust réseau"); title(s, "Aucun port ouvert : JIT SSH + Cloudflare")
# bloc SSH JIT
box(s, Inches(0.8), Inches(2.2), Inches(5.7), Inches(2.0),
    "SSH juste-à-temps",
    "Port 22 FERMÉ par défaut. La CI l'ouvre pour la seule IP du runner le temps\nd'Ansible, puis le referme — même en cas d'échec. Jamais exposé en continu.",
    fill=SOFT, accent=BRAND)
# bloc tunnel
box(s, Inches(6.8), Inches(2.2), Inches(5.7), Inches(2.0),
    "Cloudflare Tunnel + Access",
    "Les UIs (ArgoCD, Grafana, Kiali) sortent par un tunnel sortant — aucun port\nentrant. Cloudflare Access met un login OTP devant. La boutique reste publique.",
    fill=SOFT, accent=V2)
textbox(s, Inches(0.7), Inches(4.5), Inches(12), Inches(1.3),
        "Le zero-trust ne s'arrête pas aux services : l'accès à la plateforme suit le même principe.\n"
        "argocd / grafana / kiali → login obligatoire.   shop → public.   80/443/22 → fermés au repos.",
        size=15, color=INK, spacing=1.3)
textbox(s, Inches(0.7), Inches(6.2), Inches(12), Inches(0.5),
        "Connexion SORTANTE uniquement : rien n'entre depuis Internet sans passer par Cloudflare.",
        size=14, color=DIM, align=PP_ALIGN.CENTER)
footer(s, 14)

# =====================================================================
# 15 — La plateforme (infra reproductible)
# =====================================================================
s = slide(); kicker(s, "Plateforme"); title(s, "Infra reproductible : Terraform → Ansible")
steps = [("Terraform", "serveur Hetzner", 0.8),
         ("Ansible", "k3s · Istio · ArgoCD", 4.0),
         ("k3s", "cluster", 7.2),
         ("ArgoCD", "GitOps", 10.0)]
for name, sub, x in steps:
    box(s, Inches(x), Inches(3.0), Inches(2.7), Inches(1.2), name, sub, accent=BRAND)
for x in (3.5, 6.7, 9.5):
    arrow(s, Inches(x), Inches(3.6), Inches(x+0.5), Inches(3.6), color=DIM)
textbox(s, Inches(0.7), Inches(4.8), Inches(12), Inches(1.2),
        "Tout se recrée d'une commande : Terraform crée le serveur (state dans Terraform Cloud),\n"
        "Ansible installe k3s + observabilité + Istio/Kiali + ArgoCD. Validé from scratch :\n"
        "le déploiement complet remonte sans aucune étape manuelle.",
        size=14, color=DIM, spacing=1.25)
footer(s, 15)

# =====================================================================
# 16 — CI/CD GitOps
# =====================================================================
s = slide(); kicker(s, "Livraison"); title(s, "Pipeline sécurisée + GitOps versionné")
ci = ["push", "lint · SAST", "Trivy (CVE)", "Cosign", "GHCR @sha", "deploy/prod", "ArgoCD sync"]
x = 0.55
for i, step in enumerate(ci):
    acc = OK if step in ("Trivy (CVE)", "Cosign") else BRAND
    box(s, Inches(x), Inches(2.7), Inches(1.62), Inches(1.0), step, fill=SOFT, accent=acc)
    if i < len(ci)-1:
        arrow(s, Inches(x+1.62), Inches(3.2), Inches(x+1.75), Inches(3.2), color=DIM)
    x += 1.78
box(s, Inches(0.8), Inches(4.4), Inches(11.7), Inches(1.0),
    "ArgoCD surveille deploy/prod, pas main.",
    "Les manifests y portent l'image épinglée au SHA du commit. Git = source de vérité, rollback = revert.",
    fill=SOFT, accent=OK)
textbox(s, Inches(0.7), Inches(5.7), Inches(12), Inches(0.6),
        "Une CVE haute non corrigée casse le build. Pas de tag latest en prod.",
        size=14, color=DIM, align=PP_ALIGN.CENTER)
footer(s, 16)

# =====================================================================
# 17 — Récap / clôture
# =====================================================================
s = slide()
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.2), Inches(0.22), Inches(3.1))
band.fill.solid(); band.fill.fore_color.rgb = V2; band.line.fill.background(); band.shadow.inherit = False
textbox(s, Inches(0.9), Inches(2.2), Inches(11), Inches(0.6),
        "En une phrase", size=14, color=V2, bold=True)
textbox(s, Inches(0.9), Inches(2.75), Inches(11.8), Inches(2.2),
        "Le même zero-trust, d'abord câblé à la main (SPIRE),\npuis confié à un service mesh (Istio) :\nidentités vérifiées en mTLS, autorisées en déclaratif,\ndéployées en GitOps et observables.",
        size=25, color=INK, bold=True, spacing=1.15)
textbox(s, Inches(0.9), Inches(5.5), Inches(11.6), Inches(0.5),
        "SPIFFE/SPIRE → Istio/Kiali · Go + Node.js · Terraform/Ansible · ArgoCD · Prometheus/Grafana · Cloudflare",
        size=14, color=DIM)

prs.save("slides/SPIRE-Zero-Trust.pptx")
print(f"OK — {len(prs.slides.__iter__.__self__._sldIdLst)} slides")
